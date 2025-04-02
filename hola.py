from pptx import Presentation
import re

ruta_archivo = "presentacion.pptx"
ruta_salida = "presentacion_actualizada.pptx"

# Cargar presentación
prs = Presentation(ruta_archivo)

# Regex para capturar precios como $79.900 o 79.900
precio_regex = re.compile(r'\$?\d{1,3}(?:[.,]\d{3})+')

# Función para convertir texto a número
def texto_a_numero(texto):
    texto = texto.replace('$', '').replace('.', '').replace(',', '')
    try:
        return int(texto)
    except:
        return None

# Para mostrar cambios
precios_actualizados = []

# Recorremos todas las diapositivas
for slide in prs.slides:
    for shape in slide.shapes:
        if shape.has_text_frame:
            text_frame = shape.text_frame

            for paragraph in text_frame.paragraphs:
                # Combinar todos los runs en uno solo
                texto_original = ''.join(run.text for run in paragraph.runs)

                # Buscar precios
                matches = list(precio_regex.finditer(texto_original))
                if not matches:
                    continue

                nuevo_texto = texto_original
                for match in reversed(matches):  # Reemplazar de atrás hacia adelante
                    precio = match.group()
                    numero = texto_a_numero(precio)
                    if numero is not None:
                        nuevo_valor = numero + 19900
                        nuevo_precio = f"${nuevo_valor:,.0f}".replace(",", ".")
                        precios_actualizados.append((precio, nuevo_precio))
                        nuevo_texto = nuevo_texto[:match.start()] + nuevo_precio + nuevo_texto[match.end():]

                # Limpiar los runs existentes y asignar el nuevo texto completo
                for run in paragraph.runs:
                    run.text = ''
                paragraph.runs[0].text = nuevo_texto

# Guardar
prs.save(ruta_salida)
print("✅ Archivo guardado como:", ruta_salida)

# Mostrar cambios realizados
print("\n💰 Precios modificados:")
for original, nuevo in precios_actualizados:
    print(f"   {original} → {nuevo}")
