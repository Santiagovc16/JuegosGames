from pptx import Presentation

# Cargar archivo
ruta_archivo = "presentacion.pptx"
prs = Presentation(ruta_archivo)

print("\n📄 Contenido de las diapositivas:\n")

for i, slide in enumerate(prs.slides, start=1):
    print(f"--- Diapositiva {i} ---")
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                text = ''.join([run.text for run in paragraph.runs])
                print(text)
